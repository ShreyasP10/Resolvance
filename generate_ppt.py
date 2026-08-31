from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors per sih py ppt.txt
HDR = RGBColor(0x1F,0x4E,0x79)
GREEN = RGBColor(0x15,0x80,0x3D)
RED = RGBColor(0xA3,0x2B,0x20)
CREAM = RGBColor(0xFF,0xF2,0xCC)
CREAM2 = RGBColor(0xDE,0xEB,0xF7)
GREY = RGBColor(0xF2,0xF2,0xF2)
BLACK = RGBColor(0x00,0x00,0x00)
WHITE = RGBColor(0xFF,0xFF,0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
prs.slide_layouts[6]  # blank

def hex_to_rgb(h): h=h.lstrip('#'); return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, line_color=BLACK, line_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = line_width
    return shape

def add_text(slide, left, top, width, height, text, font_size=10, bold=False, color=BLACK, font_name="Calibri", alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullets(slide, left, top, width, height, bullets, font_size=9, color=BLACK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i>0 else tf.paragraphs[0]
        p.text = "• " + b
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(2)
        p.level = 0
    return txBox

# Slide 1 Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.6), HDR)
add_text(slide, Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.3), "SMART INDIA HACKATHON 2026", 10, False, WHITE, "Calibri", PP_ALIGN.CENTER)
add_text(slide, Inches(0.3), Inches(1.2), Inches(12.7), Inches(0.6), "Resolvance — Sovereign Super-Resolution for Sentinel-2", 23, True, HDR, "Times New Roman", PP_ALIGN.CENTER)
add_text(slide, Inches(0.3), Inches(1.9), Inches(12.7), Inches(0.3), "10m → 3.3m N-Channel  |  Spectral Truth (SAM)  |  Uncertainty", 11, False, BLACK, "Calibri", PP_ALIGN.CENTER)
# details boxes
add_shape(slide, Inches(0.3), Inches(2.6), Inches(12.73), Inches(2.2), WHITE)
add_text(slide, Inches(0.5), Inches(2.7), Inches(12.3), Inches(0.3), "Problem Statement ID: SIH26142  |  Title: Deep Learning Based Super Resolution Mapping 10m Sentinel-2 → <4m  |  Theme: Space Technology  |  Category: Software", 9, False, BLACK)
add_text(slide, Inches(0.5), Inches(3.1), Inches(12.3), Inches(0.3), "Organization: NTRO  |  Team: Antariksh Setu (Resolvance)  |  PS Scale: 10m → 3.3m (3×)", 9, False, BLACK)
add_text(slide, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.4), "Project: Fresh Python+PyTorch pipeline — N-band GeoTIFF in → COG GeoTIFF + heatmap out — sovereign, on-prem, QGIS-ready", 9, False, BLACK)
add_text(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.3), "KPI: PSNR≥28 SSIM≥0.80 SAM<3° NDVI r>0.98 RMSE<0.3px  |  Demo ≤30s CPU  |  1024×1024×4", 8, False, BLACK, "Consolas")
add_text(slide, Inches(0.3), Inches(6.8), Inches(12.7), Inches(0.3), "Team Antariksh Setu  |  SIH26142  |  Space Tech", 8, False, BLACK, "Calibri", PP_ALIGN.CENTER)

# Slide 2 Idea
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.5), HDR)
add_text(slide, Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.3), "IDEA TITLE & PROPOSED SOLUTION", 14, True, WHITE, "Calibri", PP_ALIGN.CENTER)
add_text(slide, Inches(0.3), Inches(0.6), Inches(12.7), Inches(0.3), "Resolvance — Sovereign Super-Resolution for Sentinel-2 10m→3.3m N-Channel with Spectral Truth & Uncertainty  |  Antariksh Setu", 9, True, BLACK, "Calibri", PP_ALIGN.CENTER)
# Problem
add_shape(slide, Inches(0.3), Inches(1.0), Inches(6.4), Inches(1.0), WHITE)
add_text(slide, Inches(0.4), Inches(1.05), Inches(6.2), Inches(0.2), "Problem:", 9, True, RED)
add_text(slide, Inches(0.4), Inches(1.25), Inches(6.2), Inches(0.7), "Free Sentinel-2 covers India every 5 days at 10m, but 10m hides roads, buildings and field edges — and <1m commercial data costs $1000s. Without 10m→<4m SR, analysts work blind or overspend.", 8, False, BLACK)
# Idea
add_shape(slide, Inches(0.3), Inches(2.1), Inches(6.4), Inches(1.1), CREAM)
add_text(slide, Inches(0.4), Inches(2.15), Inches(6.2), Inches(0.2), "Our Idea:", 9, True, HDR)
add_text(slide, Inches(0.4), Inches(2.35), Inches(6.2), Inches(0.8), "Resolvance is an AI-powered pipeline with a three-layer approach — N-channel SR, spectral SAM guard, and uncertainty heatmap — ensuring reliable <4m mapping when 10m falls short. It uses tiled ESRGAN/SwinIR 3× trained on synthetic SpaceNet pairs to deliver sharp, spectral-true, trust-scored COGs.", 7, False, BLACK)
# How addresses
add_shape(slide, Inches(0.3), Inches(3.3), Inches(6.4), Inches(1.2), WHITE)
add_text(slide, Inches(0.4), Inches(3.35), Inches(6.2), Inches(0.2), "How It Addresses Problem:", 8, True, GREEN)
add_bullets(slide, Inches(0.4), Inches(3.55), Inches(6.2), Inches(0.9), ["Eliminates $1000s commercial buy — free 5-day revisit → <4m", "On-prem 2.6s CPU (1024²×4 tiled), no foreign API", "Tactical visibility on narrow roads/buildings/fields — QGIS COG EPSG:32643"], 7)
# Right innovation 4 boxes
add_text(slide, Inches(6.95), Inches(1.0), Inches(6.05), Inches(0.2), "INNOVATION & UNIQUENESS", 9, True, HDR, "Calibri", PP_ALIGN.CENTER)
boxes = [
    ("Self-Supervised", "SpaceNet 0.5m→10m synthetic pairs solve paired-data gap prompt.txt:103"),
    ("Spectral Fidelity", "N-channel (4/8-band) + SAM loss — no RGB hallucination SAM<3°"),
    ("Confidence Maps", "MC-Dropout T=10 viridis — flags hallucinations for defence >0.6 red"),
    ("Vendor-Free", "ONNX + COG — swap ESRGAN/SwinIR, QGIS-ready, no lock-in"),
]
for i, (t, d) in enumerate(boxes):
    y = Inches(1.25 + i*0.85)
    add_shape(slide, Inches(6.95), y, Inches(6.05), Inches(0.75), CREAM if i%2==0 else CREAM2)
    add_text(slide, Inches(7.05), y+Inches(0.05), Inches(5.85), Inches(0.2), t, 8, True, HDR)
    add_text(slide, Inches(7.05), y+Inches(0.25), Inches(5.85), Inches(0.4), d, 7, False, BLACK)

# Slide 3 Technical Approach
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.5), HDR)
add_text(slide, Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.3), "TECHNICAL APPROACH", 14, True, WHITE, "Calibri", PP_ALIGN.CENTER)
# Left stack 4 boxes
left_items = [
    ("PIPELINE & IO", "Python Flask + Gunicorn\nPOST /api/infer ≤50MB\ntifffile + Affine C×H×W 33550/33922/34735"),
    ("INFERENCE", "PyTorch ESRGAN/SwinIR 1×1 N→64\nscale=3 10m→3.3m 256→768\nStub bicubic → weights/*.pth"),
    ("GEOSPATIAL & DEMO", "TileManager 256×16 Gaussian 0.1-1.0\nacc/wsum Affine.scale\nVanilla JS + Leaflet COG EPSG:32643"),
    ("TRAINING, LOSS & METRICS", "SpaceNet 0.5m→10m blur+INTER_AREA\nL1 +0.1·Perceptual+0.05·SAM\nPSNR/SSIM/SAM<3° NDVI>0.98 RMSE<0.3px"),
]
for i, (t, d) in enumerate(left_items):
    y = Inches(0.7 + i*1.15)
    add_shape(slide, Inches(0.3), y, Inches(6.4), Inches(1.05), CREAM if i%2==0 else WHITE)
    add_text(slide, Inches(0.4), y+Inches(0.05), Inches(6.2), Inches(0.15), t, 7, True, HDR)
    add_text(slide, Inches(0.4), y+Inches(0.22), Inches(6.2), Inches(0.8), d, 6, False, BLACK, "Consolas")
# Right flow
add_shape(slide, Inches(6.95), Inches(0.7), Inches(6.05), Inches(2.8), WHITE)
add_text(slide, Inches(7.05), Inches(0.75), Inches(5.85), Inches(0.2), "PIPELINE IMPLEMENTATION & EXPORT FLOW", 7, True, HDR, "Calibri", PP_ALIGN.CENTER)
flow = ["INPUT 10m 4-band EPSG:32643", "PATCH TILE 256×16 Windowed", "N-CH SR 3× ESRGAN/SwinIR", "SAM + UNCERTAINTY HEAD", "STITCH Gaussian Blend", "EXPORT COG + Heatmap", "QGIS Validation"]
for i, f in enumerate(flow):
    y = Inches(1.0 + i*0.33)
    bg = CREAM2 if "HEAT" in f or "SAM" in f else (GREY if "TILE" in f else WHITE)
    add_shape(slide, Inches(7.2), y, Inches(5.5), Inches(0.28), bg)
    add_text(slide, Inches(7.3), y+Inches(0.02), Inches(5.3), Inches(0.24), f, 6, False, BLACK, "Calibri", PP_ALIGN.CENTER)
    if i < len(flow)-1:
        add_text(slide, Inches(9.7), y+Inches(0.28), Inches(0.5), Inches(0.1), "↓", 6, True, HDR, "Calibri", PP_ALIGN.CENTER)
add_text(slide, Inches(7.05), Inches(3.5), Inches(5.85), Inches(0.3), "Synthetic Loop: SpaceNet 0.5m → Blur/Downsample → synthetic 10m", 6, True, GREEN, "Calibri", PP_ALIGN.CENTER)
# Bottom methodology
add_shape(slide, Inches(0.3), Inches(5.4), Inches(12.73), Inches(1.1), GREY)
add_text(slide, Inches(0.4), Inches(5.45), Inches(12.5), Inches(0.15), "METHODOLOGY", 7, True, HDR, "Calibri", PP_ALIGN.CENTER)
steps = ["Read\nN-band GeoTIFF", "Tile\nWindowed 256+16", "SR\nN-ch SwinIR", "SAM+Heat\nLoss+Dropout", "Stitch\nGaussian Blend", "Validate+Export\nCOG + QGIS"]
for i, s in enumerate(steps):
    x = Inches(0.5 + i*2.05)
    add_shape(slide, Inches(0.5 + i*2.05), Inches(5.65), Inches(1.85), Inches(0.65), WHITE)
    add_text(slide, x+Inches(0.05), Inches(5.68), Inches(1.75), Inches(0.6), s, 6, False, BLACK, "Calibri", PP_ALIGN.CENTER)
    if i < 5:
        add_text(slide, x+Inches(1.85), Inches(5.9), Inches(0.2), Inches(0.2), "→", 8, True, HDR, "Calibri", PP_ALIGN.CENTER)

# Slide 4 Feasibility
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.5), HDR)
add_text(slide, Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.3), "FEASIBILITY AND VIABILITY", 14, True, WHITE, "Calibri", PP_ALIGN.CENTER)
add_text(slide, Inches(0.3), Inches(0.5), Inches(12.7), Inches(0.15), "Antariksh Setu  |  SMART INDIA HACKATHON 2026", 7, False, BLACK, "Calibri", PP_ALIGN.CENTER)
# Left feasibility
add_text(slide, Inches(0.3), Inches(0.75), Inches(6.4), Inches(0.2), "Analysis of Feasibility", 8, True, HDR)
feas = [
    ("Technical", "Stateless Flask + tiled ML <8GB, 1024² CPU <30s, GPU <10s"),
    ("Operational", "Single Docker air-gapped for NTRO on-prem, no dependencies"),
    ("Economic", "Open Sentinel-2+SpaceNet, Colab Pro ₹2k vs Maxar ₹5-10L/10k km²"),
]
for i, (t,d) in enumerate(feas):
    y = Inches(0.95 + i*0.55)
    add_shape(slide, Inches(0.3), y, Inches(6.4), Inches(0.45), CREAM if i%2==0 else WHITE)
    add_text(slide, Inches(0.4), y+Inches(0.02), Inches(6.2), Inches(0.15), t, 7, True, GREEN)
    add_text(slide, Inches(0.4), y+Inches(0.18), Inches(6.2), Inches(0.25), d, 6, False, BLACK)
# Middle risks
add_text(slide, Inches(0.3), Inches(2.7), Inches(12.73), Inches(0.2), "Potential Challenges and Risks", 8, True, RED, "Calibri", PP_ALIGN.CENTER)
risks = ["Data Pairs: No real same-day 10m↔3m", "Hallucination: GAN fakes NDVI", "OOM: 5000²×4 → 1GB", "8-band Crash: on-spot prompt.txt:98"]
for i, r in enumerate(risks):
    x = Inches(0.3 + i*3.2)
    add_shape(slide, Inches(0.3 + i*3.2), Inches(2.95), Inches(3.0), Inches(0.45), GREY)
    add_text(slide, x+Inches(0.05), Inches(3.02), Inches(2.9), Inches(0.35), r, 6, False, BLACK, "Calibri", PP_ALIGN.CENTER)
# Right strategies
add_text(slide, Inches(0.3), Inches(3.5), Inches(12.73), Inches(0.2), "Strategies for Overcoming Challenges", 8, True, GREEN, "Calibri", PP_ALIGN.CENTER)
strats = [
    ("Synthetic Degradation", "SpaceNet 0.5m→blur σ0.8+INTER_AREA 6×+1% noise"),
    ("Uncertainty Heatmap", "MC-Dropout T=10 viridis >0.6 red manual review"),
    ("Dynamic Tiling", "256+16 stride240 Gaussian acc/wsum tifffile"),
    ("N-Channel Factory", "1×1 conv mean-init auto N=1/3/4/8 no crash"),
]
for i, (t,d) in enumerate(strats):
    x = Inches(0.3 + i*3.2)
    add_shape(slide, Inches(0.3 + i*3.2), Inches(3.75), Inches(3.0), Inches(0.7), CREAM if i%2==0 else CREAM2)
    add_text(slide, x+Inches(0.05), Inches(3.8), Inches(2.9), Inches(0.15), t, 7, True, HDR)
    add_text(slide, x+Inches(0.05), Inches(4.0), Inches(2.9), Inches(0.4), d, 6, False, BLACK)

# Slide 5 Impact
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.5), HDR)
add_text(slide, Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.3), "IMPACT AND BENEFITS", 14, True, WHITE, "Calibri", PP_ALIGN.CENTER)
add_text(slide, Inches(0.3), Inches(0.5), Inches(12.7), Inches(0.15), "Antariksh Setu  |  SMART INDIA HACKATHON 2026", 7, False, BLACK, "Calibri", PP_ALIGN.CENTER)
add_text(slide, Inches(0.3), Inches(0.7), Inches(12.7), Inches(0.2), "Potential Impact on Target Audience", 8, True, HDR, "Calibri", PP_ALIGN.CENTER)
impacts = [
    ("Defence & Intelligence (NTRO)", "Sovereign sub-4m tactical visibility; eliminates foreign Maxar buys."),
    ("Government & Disaster (ISRO, NDMA)", "Panchayat-level flood/urban mapping without new surveys."),
    ("AgTech & Insurance", "Field boundary + post-disaster crop damage without drones."),
    ("Analyst Trust", "Uncertainty heatmap 0-1 manages liability from hallucinations."),
]
for i, (t,d) in enumerate(impacts):
    x = Inches(0.3 + i*3.2)
    add_shape(slide, Inches(0.3 + i*3.2), Inches(0.95), Inches(3.0), Inches(1.0), CREAM if i%2==0 else WHITE)
    add_text(slide, x+Inches(0.05), Inches(1.0), Inches(2.9), Inches(0.25), t, 7, True, HDR)
    add_text(slide, x+Inches(0.05), Inches(1.3), Inches(2.9), Inches(0.6), d, 6, False, BLACK)
add_text(slide, Inches(0.3), Inches(2.1), Inches(12.7), Inches(0.2), "Benefits of the Solution", 8, True, HDR, "Calibri", PP_ALIGN.CENTER)
benefits = [
    ("Economic Value 📊", "₹111Cr/year @₹5/km² saved; reduces licenses"),
    ("Strategic & Sovereign 🛡️", "On-prem air-gapped; secures national EO pipeline"),
    ("Environmental & Scientific 🍃", "Preserves NIR+NDVI for climate tracking"),
]
for i, (t,d) in enumerate(benefits):
    x = Inches(0.3 + i*4.2)
    add_shape(slide, Inches(0.3 + i*4.2), Inches(2.35), Inches(4.0), Inches(0.9), GREY if i==1 else CREAM)
    add_text(slide, x+Inches(0.05), Inches(2.4), Inches(3.9), Inches(0.2), t, 8, True, HDR, "Calibri", PP_ALIGN.CENTER)
    add_text(slide, x+Inches(0.05), Inches(2.7), Inches(3.9), Inches(0.5), d, 7, False, BLACK)

# Slide 6 Research
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.5), HDR)
add_text(slide, Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.3), "RESEARCH AND REFERENCES", 14, True, WHITE, "Calibri", PP_ALIGN.CENTER)
add_text(slide, Inches(0.3), Inches(0.5), Inches(12.7), Inches(0.15), "Antariksh Setu  |  SMART INDIA HACKATHON 2026", 7, False, BLACK, "Calibri", PP_ALIGN.CENTER)
add_shape(slide, Inches(0.3), Inches(0.7), Inches(4.0), Inches(5.5), CREAM)
add_text(slide, Inches(0.35), Inches(0.75), Inches(3.9), Inches(0.2), "Key Research Papers", 8, True, HDR)
add_bullets(slide, Inches(0.35), Inches(0.95), Inches(3.9), Inches(2.0), [
    "Real-ESRGAN: Wang et al., ICCVW 2021 DOI 10.1109/ICCVW54120",
    "SwinIR: Liang et al., ICCVW 2021",
    "DSen2: Lanaras et al., RSE 2018 DOI 10.1016/j.rse",
    "MC-Dropout: Gal & Ghahramani, ICML 2016 PMLR 48",
    "SAM: Kruse et al., RSE 1993 DOI 10.1016/0034",
], 6)
add_text(slide, Inches(0.35), Inches(3.0), Inches(3.9), Inches(0.2), "Open Datasets & Standards", 8, True, HDR)
add_bullets(slide, Inches(0.35), Inches(3.2), Inches(3.9), Inches(1.5), [
    "SpaceNet 0.5m spacenet.ai prompt.txt:103",
    "Sentinel-2 L2A 10m Copernicus",
    "COG Standard OGC tifffile 33550/33922/34735",
    "SIH2026_RESEARCH: no sovereign SR",
], 6)
add_shape(slide, Inches(4.6), Inches(0.7), Inches(4.0), Inches(5.5), CREAM2)
add_text(slide, Inches(4.65), Inches(0.75), Inches(3.9), Inches(0.2), "Core Frameworks", 8, True, HDR)
add_bullets(slide, Inches(4.65), Inches(0.95), Inches(3.9), Inches(2.0), [
    "Rasterio/GDAL + Affine core/io.py",
    "PyTorch + torchvision Real-ESRGAN core/transforms.py",
    "GeoTIFF.js + Leaflet frontend/static/js/app.js",
    "scikit-image / lpips / albumentations core/metrics.py",
], 6)
add_shape(slide, Inches(8.9), Inches(0.7), Inches(4.13), Inches(5.5), WHITE)
add_text(slide, Inches(8.95), Inches(0.75), Inches(4.0), Inches(0.2), "Traceability", 8, True, GREEN)
add_bullets(slide, Inches(8.95), Inches(0.95), Inches(4.0), Inches(2.0), [
    "prompt.txt:47-52 Expected Solution → PRD 7",
    "prompt.txt:64 SAM + NDVI → SRS 3.6",
    "prompt.txt:98 8-band on-spot → PR-02",
    "prompt.txt:106-115 Judge Q&A → DRD/FAQ",
    "Kepler-404 → Resolvance fresh v2:6",
], 6)

out = "Resolvance_SIH26142_AntarikshSetu.pptx"
prs.save(out)
print(f"saved {out}")
